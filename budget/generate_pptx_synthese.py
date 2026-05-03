"""
Génère un PowerPoint de synthèse budget artistes - 1 diapo.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

# Background
bg = slide.background
fill = bg.fill
fill.solid()
fill.fore_color.rgb = RGBColor(255, 255, 255)

# Title
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(0.8))
tf = title_box.text_frame
p = tf.paragraphs[0]
p.text = "MONIQUE FESTIVAL 2026 — Budget Artistes"
p.font.size = Pt(28)
p.font.bold = True
p.font.color.rgb = RGBColor(0, 0, 0)
p.alignment = PP_ALIGN.CENTER

# Subtitle
sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.9), Inches(12), Inches(0.4))
tf2 = sub_box.text_frame
p2 = tf2.paragraphs[0]
p2.text = "Synthèse au 04/04/2026 — basé sur le Line Up final"
p2.font.size = Pt(14)
p2.font.color.rgb = RGBColor(100, 100, 100)
p2.alignment = PP_ALIGN.CENTER

# --- LEFT COLUMN: Règles ---
rules_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(5.5), Inches(5.5))
tf_rules = rules_box.text_frame
tf_rules.word_wrap = True

def add_rule(tf, text, bold_part="", size=12):
    p = tf.add_paragraph()
    if bold_part:
        run = p.add_run()
        run.text = bold_part
        run.font.bold = True
        run.font.size = Pt(size)
        run.font.color.rgb = RGBColor(0, 0, 0)
        run2 = p.add_run()
        run2.text = text
        run2.font.size = Pt(size)
        run2.font.color.rgb = RGBColor(50, 50, 50)
    else:
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.color.rgb = RGBColor(50, 50, 50)
    p.space_after = Pt(4)

# Section title
p_title = tf_rules.paragraphs[0]
p_title.text = "RÈGLES DE CHIFFRAGE"
p_title.font.size = Pt(16)
p_title.font.bold = True
p_title.font.color.rgb = RGBColor(220, 50, 50)
p_title.space_after = Pt(10)

add_rule(tf_rules, " 250 € coût chargé / prestation (153 € brut min. x1,6)", "1. Cachet spectacle :")
add_rule(tf_rules, " 100 € coût chargé / service (62 € brut / 3h)\n     Ateliers = répétitions en vue de la Création Monique", "2. Répétition atelier :")
add_rule(tf_rules, " 4 acteurs estimés / créneau\n     Artistes locaux → pas de défraiement", "3. Théâtre appel à projet :")
add_rule(tf_rules, " 0 € (Léopold, Lorraine, Luther, Lewis)", "4. Famille :")
add_rule(tf_rules, " 200 € max / artiste non-famille\n     Prise en charge réelle estimée à 75% (-25%)", "5. Défraiement :")
add_rule(tf_rules, " prix coûtant (repas 3 €, boisson 1,50 €)", "6. Restauration :")

# --- Effectifs ---
eff_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(5.5), Inches(1.5))
tf_eff = eff_box.text_frame
tf_eff.word_wrap = True
p_eff = tf_eff.paragraphs[0]
p_eff.text = "EFFECTIFS"
p_eff.font.size = Pt(16)
p_eff.font.bold = True
p_eff.font.color.rgb = RGBColor(220, 50, 50)
p_eff.space_after = Pt(8)

for line in [
    "55 artistes (dont 4 famille)",
    "58 cachets spectacle payés",
    "12 répétitions atelier payées",
    "32 artistes défrayés",
]:
    p = tf_eff.add_paragraph()
    p.text = "• " + line
    p.font.size = Pt(12)
    p.space_after = Pt(2)

# --- RIGHT COLUMN: Budget ---
from pptx.util import Inches, Pt

# Budget table
rows, cols = 8, 3
tbl_shape = slide.shapes.add_table(rows, cols, Inches(6.5), Inches(1.5), Inches(6.3), Inches(3.5))
tbl = tbl_shape.table

# Column widths
tbl.columns[0].width = Inches(3.2)
tbl.columns[1].width = Inches(1.5)
tbl.columns[2].width = Inches(1.6)

data = [
    ("Poste", "Détail", "Total"),
    ("Cachets spectacle", "58 x 250 €", "14 500 €"),
    ("Répétitions atelier", "12 x 100 €", "1 200 €"),
    ("Sous-total prestations", "", "15 700 €"),
    ("Défraiement (max 200 €)", "32 artistes", "6 400 €"),
    ("Ajustement -25%", "prise en charge réelle", "-1 600 €"),
    ("Sous-total défraiement", "", "4 800 €"),
    ("TOTAL ARTISTES", "", "20 500 €"),
]

for i, (c1, c2, c3) in enumerate(data):
    for j, val in enumerate([c1, c2, c3]):
        cell = tbl.cell(i, j)
        cell.text = val
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(11)
            if i == 0:
                paragraph.font.bold = True
                paragraph.font.color.rgb = RGBColor(255, 255, 255)
            elif i in (3, 6):
                paragraph.font.bold = True
                paragraph.font.color.rgb = RGBColor(0, 0, 100)
            elif i == 7:
                paragraph.font.bold = True
                paragraph.font.size = Pt(13)
                paragraph.font.color.rgb = RGBColor(220, 50, 50)
            if j == 2:
                paragraph.alignment = PP_ALIGN.RIGHT

    # Header row color
    if i == 0:
        for j in range(3):
            tbl.cell(i, j).fill.solid()
            tbl.cell(i, j).fill.fore_color.rgb = RGBColor(50, 50, 50)
    elif i == 7:
        for j in range(3):
            tbl.cell(i, j).fill.solid()
            tbl.cell(i, j).fill.fore_color.rgb = RGBColor(255, 235, 235)
    elif i in (3, 6):
        for j in range(3):
            tbl.cell(i, j).fill.solid()
            tbl.cell(i, j).fill.fore_color.rgb = RGBColor(230, 230, 245)

# --- Par jour table ---
rows2, cols2 = 5, 6
tbl2_shape = slide.shapes.add_table(rows2, cols2, Inches(6.5), Inches(5.3), Inches(6.3), Inches(1.8))
tbl2 = tbl2_shape.table

tbl2.columns[0].width = Inches(1.2)
tbl2.columns[1].width = Inches(0.9)
tbl2.columns[2].width = Inches(1.1)
tbl2.columns[3].width = Inches(0.9)
tbl2.columns[4].width = Inches(1.1)
tbl2.columns[5].width = Inches(1.1)

data2 = [
    ("Jour", "Nb cachets", "Cachets", "Nb rép.", "Répétitions", "Total"),
    ("Vendredi", "16", "4 000 €", "0", "0 €", "4 000 €"),
    ("Samedi", "20", "5 000 €", "6", "600 €", "5 600 €"),
    ("Dimanche", "22", "5 500 €", "6", "600 €", "6 100 €"),
    ("TOTAL", "58", "14 500 €", "12", "1 200 €", "15 700 €"),
]

for i, row_data in enumerate(data2):
    for j, val in enumerate(row_data):
        cell = tbl2.cell(i, j)
        cell.text = val
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(11)
            if i == 0:
                paragraph.font.bold = True
                paragraph.font.color.rgb = RGBColor(255, 255, 255)
            elif i == 4:
                paragraph.font.bold = True
            if j > 0:
                paragraph.alignment = PP_ALIGN.RIGHT

    if i == 0:
        for j in range(6):
            tbl2.cell(i, j).fill.solid()
            tbl2.cell(i, j).fill.fore_color.rgb = RGBColor(50, 50, 50)
    elif i == 4:
        for j in range(6):
            tbl2.cell(i, j).fill.solid()
            tbl2.cell(i, j).fill.fore_color.rgb = RGBColor(230, 230, 230)

output = "c:/Users/charl/Monique-Festival/budget/Synthese_budget_artistes.pptx"
prs.save(output)
print(f"OK: {output}")
